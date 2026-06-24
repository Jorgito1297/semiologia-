#!/usr/bin/env python3
"""
Módulo de Extracción y Sanitización de Documentos (document_parser.py)
Medical AI Learning System — MED-228-T-1 (UCE)

Este módulo se encarga de:
1. Asegurar la instalación automática de dependencias críticas en la PC del estudiante.
2. Extraer texto plano estructurado de archivos PDF (soporta doble columna para libros).
3. Extraer contenido de PPTX (diapositivas, tablas y notas del orador/docente).
4. Extraer contenido de DOCX (guías, manuales de práctica y tablas).
5. Limpiar y sanitizar transcripciones WebVTT y SRT eliminando marcas de tiempo, cabeceras y muletillas.
6. Transcribir de forma multimodal grabaciones de audio/video (.mp3, .mp4) subiéndolas a Gemini File API.
"""

import os
import sys
import re
import importlib
import subprocess
import unicodedata

# =====================================================================
# ⚙️ GESTIÓN DE DEPENDENCIAS AUTOMÁTICA
# =====================================================================
def ensure_dependencies():
    """
    Verifica e instala de forma automática las dependencias requeridas
    para asegurar una ejecución zero-friction en la computadora del estudiante.
    """
    dependencies = {
        "fitz": "pymupdf",
        "pptx": "python-pptx",
        "docx": "python-docx",
        "watchdog": "watchdog",
        "requests": "requests",
        "google.genai": "google-genai"
    }
    
    missing = []
    for module_name, pip_name in dependencies.items():
        try:
            importlib.import_module(module_name)
        except ImportError:
            missing.append(pip_name)
            
    if missing:
        print(f"[PARSER]: Instalando dependencias faltantes: {', '.join(missing)}...")
        try:
            # Ejecuta pip install de manera silenciosa
            subprocess.check_call([sys.executable, "-m", "pip", "install", *missing])
            print("[PARSER]: Todas las dependencias se instalaron correctamente.")
        except Exception as e:
            print(f"[PARSER]: Error crítico al instalar dependencias automáticas: {e}")
            print("[PARSER]: Por favor ejecute: pip install " + " ".join(missing))

# Ejecutar verificación al importar o cargar el módulo
ensure_dependencies()

# Importaciones protegidas tras la instalación
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

# =====================================================================
# 📄 PARSER 1: EXTRACTOR PDF (Doble Columna & Libros Base)
# =====================================================================
def extract_pdf(pdf_path: str) -> str:
    """
    Extrae texto de un archivo PDF ordenando los bloques por columnas físicas
    de lectura, optimizado para diseños de dos columnas de libros médicos (Llanio).
    """
    if not os.path.exists(pdf_path):
        return f"Error: El archivo PDF '{pdf_path}' no existe."
        
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        return f"Error al abrir el PDF {pdf_path}: {e}"
        
    extracted_pages = []
    
    for page_num, page in enumerate(doc):
        # blocks: (x0, y0, x1, y1, "texto", block_no, block_type)
        blocks = page.get_text("blocks")
        
        page_width = page.rect.width
        mid_x = page_width / 2
        
        # Filtrar bloques de texto válidos y evitar bloques no textuales o vacíos
        valid_blocks = []
        for b in blocks:
            x0, y0, x1, y1, text, block_no, block_type = b
            if block_type == 0 and text.strip():
                valid_blocks.append(b)
                
        if not valid_blocks:
            extracted_pages.append(f"--- [PÁGINA {page_num + 1}] ---\n")
            continue
            
        # Determinar si cada bloque es cruzado (ancho completo) o específico de columna
        cross_blocks = []
        col_blocks = []
        for b in valid_blocks:
            x0, y0, x1, y1, text, block_no, block_type = b
            block_width = x1 - x0
            # Es cruzado si ocupa más del 60% de la página o si cruza claramente el eje del medio
            is_cross = (block_width > page_width * 0.6) or (x0 < mid_x - 15 and x1 > mid_x + 15)
            if is_cross:
                cross_blocks.append(b)
            else:
                col_blocks.append(b)
                
        # Clasificación global preliminar para detectar si la página tiene dos columnas principales
        col1_global = [b for b in col_blocks if (b[0] + b[2])/2 < mid_x]
        col2_global = [b for b in col_blocks if (b[0] + b[2])/2 >= mid_x]
        is_double_column = len(col1_global) > 2 and len(col2_global) > 2
        
        if is_double_column:
            # Ordenar bloques de ancho completo verticalmente
            cross_blocks.sort(key=lambda x: x[1])
            
            # Agrupar bloques de columna en bandas horizontales definidas por la posición de los bloques cruzados
            col_bands = [[] for _ in range(len(cross_blocks) + 1)]
            for b in col_blocks:
                center_y = (b[1] + b[3]) / 2
                band_idx = 0
                for idx, cb in enumerate(cross_blocks):
                    if center_y < cb[1]:
                        break
                    band_idx = idx + 1
                col_bands[band_idx].append(b)
                
            sorted_blocks = []
            for i in range(len(cross_blocks) + 1):
                # Extraer y ordenar bloques de la columna en esta banda
                b_cols = col_bands[i]
                if b_cols:
                    b_col1 = [b for b in b_cols if (b[0] + b[2])/2 < mid_x]
                    b_col2 = [b for b in b_cols if (b[0] + b[2])/2 >= mid_x]
                    b_col1.sort(key=lambda x: x[1])
                    b_col2.sort(key=lambda x: x[1])
                    sorted_blocks.extend(b_col1)
                    sorted_blocks.extend(b_col2)
                    
                # Insertar el bloque cruzado correspondiente
                if i < len(cross_blocks):
                    sorted_blocks.append(cross_blocks[i])
        else:
            # Columna única, ordenar de arriba a abajo y luego de izquierda a derecha
            valid_blocks.sort(key=lambda x: (x[1], x[0]))
            sorted_blocks = valid_blocks
            
        page_text = []
        for b in sorted_blocks:
            page_text.append(b[4].strip())
            
        page_content = "\n".join(page_text)
        extracted_pages.append(f"--- [PÁGINA {page_num + 1}] ---\n{page_content}")
        
    return "\n\n".join(extracted_pages)

# =====================================================================
# 🖥️ CONVERSORES COM (Automatización Office para formatos legacy)
# =====================================================================
def convert_ppt_to_pptx(ppt_path: str) -> str:
    """
    Convierte un archivo .ppt legacy a .pptx moderno usando PowerPoint COM.
    Retorna la ruta del archivo .pptx generado.
    """
    import win32com.client
    ppt_path_abs = os.path.abspath(ppt_path)
    pptx_path_abs = os.path.splitext(ppt_path_abs)[0] + ".pptx"
    
    if os.path.exists(pptx_path_abs):
        return pptx_path_abs
        
    print(f"[PARSER COM]: Convirtiendo {os.path.basename(ppt_path)} a PPTX...")
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # Open(FileName, ReadOnly, Untitled, WithWindow)
        presentation = powerpoint.Presentations.Open(ppt_path_abs, ReadOnly=True, Untitled=False, WithWindow=False)
        # 24 = ppSaveAsOpenXMLPresentation (formato pptx estándar)
        presentation.SaveAs(pptx_path_abs, 24)
        print(f"[PARSER COM]: Conversión exitosa -> {os.path.basename(pptx_path_abs)}")
        return pptx_path_abs
    except Exception as e:
        print(f"[PARSER COM ERROR]: Error al convertir PPT a PPTX '{ppt_path}': {e}")
        raise e
    finally:
        if presentation:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint:
            try:
                powerpoint.Quit()
            except Exception:
                pass

def convert_doc_to_docx(doc_path: str) -> str:
    """
    Convierte un archivo .doc legacy a .docx moderno usando Word COM.
    Retorna la ruta del archivo .docx generado.
    """
    import win32com.client
    doc_path_abs = os.path.abspath(doc_path)
    docx_path_abs = os.path.splitext(doc_path_abs)[0] + ".docx"
    
    if os.path.exists(docx_path_abs):
        return docx_path_abs
        
    print(f"[PARSER COM]: Convirtiendo {os.path.basename(doc_path)} a DOCX...")
    word = None
    document = None
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        document = word.Documents.Open(doc_path_abs, ReadOnly=True)
        # 12 = wdFormatXMLDocument (formato docx estándar)
        document.SaveAs2(docx_path_abs, FileFormat=12)
        print(f"[PARSER COM]: Conversión exitosa -> {os.path.basename(docx_path_abs)}")
        return docx_path_abs
    except Exception as e:
        print(f"[PARSER COM ERROR]: Error al convertir DOC a DOCX '{doc_path}': {e}")
        raise e
    finally:
        if document:
            try:
                document.Close()
            except Exception:
                pass
        if word:
            try:
                word.Quit()
            except Exception:
                pass

# =====================================================================
# 📊 PARSER 2: PRESENTACIONES PPTX (Diapositivas & Notas del Docente)
# =====================================================================
def extract_pptx(pptx_path: str) -> str:
    """
    Extrae texto de presentaciones PowerPoint conservando diapositivas,
    tablas internas (con alineación de columnas) y las notas del orador (donde el docente pone la teoría),
    soportando formas agrupadas de manera recursiva.
    """
    if not os.path.exists(pptx_path):
        return f"Error: El archivo PPTX '{pptx_path}' no existe."
        
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return f"Error al abrir la presentación PPTX {pptx_path}: {e}"
        
    def get_all_shapes(shape_container):
        flat_shapes = []
        for shape in shape_container:
            # Si es un grupo de formas, iterar recursivamente
            if hasattr(shape, "shapes") and shape.shapes:
                flat_shapes.extend(get_all_shapes(shape.shapes))
            else:
                flat_shapes.append(shape)
        return flat_shapes

    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_lines = [f"=== DIAPOSITIVA {i + 1} ==="]
        
        # Obtener todas las formas (incluidas las agrupadas)
        all_shapes = get_all_shapes(slide.shapes)
        
        # Ordenar formas por su ubicación vertical en el slide para coherencia
        shapes = sorted(
            all_shapes, 
            key=lambda s: (s.top, s.left) if hasattr(s, 'top') and s.top is not None else (0, 0)
        )
        
        for shape in shapes:
            # Cajas de texto estándar
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_lines.append(text)
                    
            # Tablas dentro de diapositivas con alineación
            if hasattr(shape, "has_table") and shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_cells = []
                    for cell in row.cells:
                        if hasattr(cell, "is_spanned") and cell.is_spanned:
                            row_cells.append("")
                        else:
                            row_cells.append(cell.text.strip().replace("\n", " "))
                    table_data.append(row_cells)
                
                if table_data:
                    max_cols = max(len(r) for r in table_data)
                    for row_cells in table_data:
                        while len(row_cells) < max_cols:
                            row_cells.append("")
                            
                    col_widths = [0] * max_cols
                    for row_cells in table_data:
                        for col_idx, cell_text in enumerate(row_cells):
                            col_widths[col_idx] = max(col_widths[col_idx], len(cell_text))
                            
                    table_text = []
                    for row_cells in table_data:
                        aligned_cells = [cell_text.ljust(col_widths[col_idx]) for col_idx, cell_text in enumerate(row_cells)]
                        table_text.append(" | ".join(aligned_cells))
                        
                    slide_lines.append("[TABLA INTERNA]:\n" + "\n".join(table_text))
                    
        # Notas del orador (Explicaciones del Docente)
        if hasattr(slide, "has_notes_slide") and slide.has_notes_slide and slide.notes_slide:
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_lines.append(f"[EXPLICACIÓN DEL DOCENTE / NOTAS]: {notes}")
                
        slides_text.append("\n".join(slide_lines))
        
    return "\n\n".join(slides_text)

# =====================================================================
# 📝 PARSER 3: DOCUMENTOS WORD DOCX (Párrafos y Tablas Criterios)
# =====================================================================
def extract_docx(docx_path: str) -> str:
    """
    Extrae texto de documentos Word respetando la jerarquía de estilos
    e incluyendo tablas estructuradas con alineación visual de celdas y control de fusiones.
    """
    if not os.path.exists(docx_path):
        return f"Error: El archivo DOCX '{docx_path}' no existe."
        
    try:
        doc = Document(docx_path)
    except Exception as e:
        return f"Error al abrir el documento DOCX {docx_path}: {e}"
        
    docx_lines = []
    
    # 1. Párrafos con indicación de estilo jerárquico
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = paragraph.style.name if paragraph.style else "Normal"
        docx_lines.append(f"[{style_name}] {text}")
        
    # 2. Tablas integradas con alineación de columnas y resolución de fusiones
    for i, table in enumerate(doc.tables):
        docx_lines.append(f"\n--- [TABLA {i + 1} EN DOCUMENTO] ---")
        
        table_data = []
        seen_cells = set()
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                # cell._tc es el elemento XML subyacente para identificar fusiones
                cell_id = cell._tc
                if cell_id in seen_cells:
                    row_cells.append("")
                else:
                    seen_cells.add(cell_id)
                    row_cells.append(cell.text.strip().replace("\n", " "))
            table_data.append(row_cells)
            
        if not table_data:
            continue
            
        # Determinar número máximo de columnas en la tabla
        max_cols = max(len(r) for r in table_data)
        
        # Rellenar celdas faltantes en filas más cortas
        for row_cells in table_data:
            while len(row_cells) < max_cols:
                row_cells.append("")
                
        # Calcular el ancho máximo de cada columna
        col_widths = [0] * max_cols
        for row_cells in table_data:
            for col_idx, cell_text in enumerate(row_cells):
                col_widths[col_idx] = max(col_widths[col_idx], len(cell_text))
                
        # Imprimir filas con alineación visual
        for row_cells in table_data:
            aligned_cells = [cell_text.ljust(col_widths[col_idx]) for col_idx, cell_text in enumerate(row_cells)]
            docx_lines.append(" | ".join(aligned_cells))
                
    return "\n".join(docx_lines)

# =====================================================================
# 🧹 PARSER 4: SANITIZADOR DE TRANSCRIPCIONES (WebVTT / SRT / Transcripts)
# =====================================================================
def sanitize_transcript(raw_content: str) -> str:
    """
    Limpia transcripciones WebVTT o SRT de Teams. Remueve marcas de tiempo,
    cabeceras, números correlativos, comentarios, metadatos y muletillas orales.
    """
    # Eliminar BOM de UTF-8 si existe para evitar que afecte el primer secuenciador
    raw_content = raw_content.lstrip('\ufeff')
    lines = raw_content.splitlines()
    cleaned_lines = []
    
    # Expresiones regulares adaptadas para WebVTT/SRT (horas opcionales y delimitadores de decimales/milisegundos)
    pattern_time = re.compile(r'(?:\d{1,2}:)?\d{2}:\d{2}[.,]\d{3}\s*-->\s*(?:\d{1,2}:)?\d{2}:\d{2}[.,]\d{3}')
    pattern_sequence = re.compile(r'^\d+$')
    pattern_headers = re.compile(r'^(WEBVTT|Kind:|Language:|NOTE|X-TIMESTAMP-MAP)', re.IGNORECASE)
    
    in_note = False
    for line in lines:
        line = line.strip()
        
        # Ignorar comentarios multilínea (NOTE) de WebVTT
        if in_note:
            if not line:
                in_note = False
            continue
            
        if line.upper().startswith("NOTE"):
            in_note = True
            continue
            
        # Ignorar vacíos, cabeceras y secuenciadores
        if not line or pattern_headers.match(line) or pattern_sequence.match(line):
            continue
            
        # Ignorar marcas de tiempo
        if pattern_time.search(line):
            continue
            
        # Normalizar locutores de Teams: <v Dr. Rivas>Texto</v> -> Dr. Rivas: Texto
        line = re.sub(r'<v\s+([^>]+)>', r'\1: ', line)
        line = re.sub(r'</v>', '', line)
        
        # Remover cualquier otra etiqueta HTML/XML residual (como marcas de tiempo inline)
        line = re.sub(r'<[^>]+>', '', line)
        
        cleaned_lines.append(line)
        
    # Concatenar en flujo de texto continuo
    text = " ".join(cleaned_lines)
    text = re.sub(r'\s+', ' ', text)
    
    # Remover ruidos/anotaciones de sonido comunes entre corchetes
    text = re.sub(r'\[(música|risas|silencio|aplausos|tos|murmullo|ininteligible)\]', '', text, flags=re.IGNORECASE)
    
    # Limpieza de muletillas y tics orales en español
    fillers = [
        r'\b(eh|estee|bueno|o\s+sea|verdad|ehmm)\b',
        r'¿?no\?',
    ]
    for filler in fillers:
        text = re.sub(filler, '', text, flags=re.IGNORECASE)
        
    return re.sub(r'\s+', ' ', text).strip()

# =====================================================================
# 🧹 PARSER 5: NORMALIZACIÓN DE TEXTO ACADÉMICO (Codificación y Separaciones)
# =====================================================================
def normalize_academic_text(text: str) -> str:
    """
    Normaliza texto Unicode (resuelve ligaduras fi/fl), une guiones silábicos
    al final de línea y remueve marcas de agua y encabezados del Llanio.
    """
    # Resuelve ligaduras tipográficas y caracteres especiales
    text = unicodedata.normalize('NFKC', text)
    
    # Resolver saltos silábicos: ej. "cardio-\nvascular" -> "cardiovascular"
    text = re.sub(r'(\w+)-\n(\w+)', r'\1\2', text)
    
    # Unir líneas rotas simples dentro del mismo párrafo
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
    
    # Mantener separación de párrafos limpios (máximo 2 saltos de línea)
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Remover marcas de copyright del libro digital
    text = re.sub(r'Editorial Ciencias Médicas.*?\n', '', text, flags=re.IGNORECASE)
    
    return text.strip()

# =====================================================================
# 🎥 PARSER 6: TRANSCRIPCIÓN MULTIMODAL API (Audio/Video crudo)
# =====================================================================
def transcribe_media_via_gemini(file_path: str, api_key: str) -> str:
    """
    Sube grabaciones (.mp4/.mp3) directamente a Gemini File API para realizar
    la transcripción y extracción conceptual sin consumir recursos locales.
    """
    try:
        from google import genai
        from google.genai import types
    except ImportError:
        return "Error: La librería oficial 'google-genai' no está disponible."
        
    if not api_key:
        return "Error: Falta la API Key de Gemini en las variables de entorno."
        
    if not os.path.exists(file_path):
        return f"Error: El archivo multimedia '{file_path}' no existe."
        
    print(f"[PARSER]: Iniciando carga de {os.path.basename(file_path)} a Gemini File API...")
    client = genai.Client(api_key=api_key)
    
    try:
        # 1. Subir archivo a la API de archivos en la nube
        media_file = client.files.upload(file=file_path)
        
        # 2. Esperar que termine el procesamiento en los servidores de Google
        import time
        print("[PARSER]: Procesando archivo en la nube", end="")
        while media_file.state.name == "PROCESSING":
            time.sleep(10)
            media_file = client.files.get(name=media_file.name)
            print(".", end="", flush=True)
        print()
        
        if media_file.state.name == "FAILED":
            raise ValueError(f"Fallo en el procesamiento de video: {media_file.error.message}")
            
        print("[PARSER]: Archivo procesado. Solicitando transcripción estructurada a gemini-2.5-flash...")
        
        # 3. Solicitar transcripción literal con modelo multimodal
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[
                media_file,
                "Eres un transcriptor médico de precisión de la UCE. Genera una transcripción literal, "
                "estructurada y detallada de esta sesión de clase. Identifica al expositor "
                "(ej. Dr. Rivas) e incluye cada término clínico discutido de forma íntegra."
            ]
        )
        
        # 4. Eliminar el archivo de la nube por privacidad y orden
        client.files.delete(name=media_file.name)
        print("[PARSER]: Archivo eliminado en la nube exitosamente.")
        return response.text
        
    except Exception as e:
        return f"Error durante la transcripción en la API: {e}"

# =====================================================================
# 🛠️ EJECUCIÓN CLI (Pruebas unitarias directas)
# =====================================================================
if __name__ == "__main__":
    print("--- MOTOR DE DEPURACIÓN DE DOCUMENTOS (MODO PRUEBA) ---")
    if len(sys.argv) < 2:
        print("Uso: python document_parser.py <ruta_archivo> [ruta_salida]")
        sys.exit(0)
        
    filepath = sys.argv[1]
    ext = filepath.split('.')[-1].lower()
    
    print(f"Procesando archivo: {filepath} ({ext})")
    parsed_text = ""
    
    if ext == 'pdf':
        parsed_text = extract_pdf(filepath)
    elif ext == 'pptx':
        parsed_text = extract_pptx(filepath)
    elif ext == 'ppt':
        try:
            pptx_path = convert_ppt_to_pptx(filepath)
            parsed_text = extract_pptx(pptx_path)
        except Exception as e:
            parsed_text = f"Error al procesar PPT: {e}"
    elif ext == 'docx':
        parsed_text = extract_docx(filepath)
    elif ext == 'doc':
        try:
            docx_path = convert_doc_to_docx(filepath)
            parsed_text = extract_docx(docx_path)
        except Exception as e:
            parsed_text = f"Error al procesar DOC: {e}"
    elif ext in ['vtt', 'srt']:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            parsed_text = sanitize_transcript(f.read())
    elif ext in ['mp4', 'mp3', 'm4a', 'wav']:
        api_key = os.environ.get("GEMINI_API_KEY")
        if not api_key:
            # Intentar cargar de .env local
            if os.path.exists(".env"):
                with open(".env", "r") as f:
                    for line in f:
                        if line.startswith("GEMINI_API_KEY="):
                            api_key = line.split("=")[1].strip()
                            break
        if api_key:
            parsed_text = transcribe_media_via_gemini(filepath, api_key)
        else:
            parsed_text = "Error: GEMINI_API_KEY no encontrada para procesar multimedia."
    else:
        print("Error: Formato no soportado.")
        sys.exit(1)
        
    # Normalizar si hay texto extraído
    if parsed_text and not parsed_text.startswith("Error"):
        parsed_text = normalize_academic_text(parsed_text)
        
    if len(sys.argv) > 2:
        dest = sys.argv[2]
        with open(dest, "w", encoding="utf-8") as f:
            f.write(parsed_text)
        print(f"Éxito. Contenido guardado en {dest}")
    else:
        print("\n=== TEXTO EXTRAÍDO Y SANITIZADO (MUESTRA 500 CARACTERES) ===")
        print(parsed_text[:500] + "\n...")
